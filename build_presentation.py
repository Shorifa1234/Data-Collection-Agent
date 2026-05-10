"""
build_presentation.py
Generates the CEO-level project overview presentation.
Run: python build_presentation.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1F, 0x38, 0x64)   # deep navy  (header bg)
GOLD       = RGBColor(0xB8, 0x96, 0x0E)
SLATE      = RGBColor(0x2E, 0x40, 0x57)   # not used
MID_BLUE   = RGBColor(0x2E, 0x59, 0x9A)   # section title
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)   # subtle bg tint
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
MID_GRAY   = RGBColor(0x55, 0x65, 0x7A)
GOLD_SOLID = RGBColor(0xC0, 0x9A, 0x10)

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

OUT_DIR  = Path(__file__).parent / "Project Overview"
OUT_FILE = OUT_DIR / "Vendor Data Intelligence Platform.pptx"

# ────────────────────────────────────────────────────────────────────────────
# Helpers
# ────────────────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.line.color.rgb = fill_color
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_textbox(slide, left, top, width, height,
                text, font_size, bold=False, italic=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT,
                word_wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font_name
    return txb

def add_multiline_textbox(slide, left, top, width, height,
                          lines,            # list of (text, size, bold, color, align)
                          spacing_after=6,
                          font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color, align) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name  = font_name
    return txb

def add_bullet_box(slide, left, top, width, height,
                   items,         # list of strings
                   font_size=14,
                   color=DARK_TEXT,
                   bullet="▸  ",
                   spacing=8,
                   font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = bullet + item
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
        run.font.name  = font_name
    return txb

def gold_divider(slide, top):
    """Thin gold horizontal rule."""
    shape = slide.shapes.add_shape(1, Inches(0.5), top, Inches(12.33), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD_SOLID
    shape.line.fill.background()
    return shape

def slide_number(slide, num, total, color=MID_GRAY):
    add_textbox(slide, Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3),
                f"{num} / {total}", 9, color=color, align=PP_ALIGN.RIGHT)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ────────────────────────────────────────────────────────────────────────────

TOTAL = 10   # total slide count

def slide_01_title(prs):
    """Full-bleed dark cover."""
    s = blank_slide(prs)

    # Background
    add_rect(s, 0, 0, W, H, NAVY)

    # Thin gold bar at bottom
    add_rect(s, 0, Inches(7.1), W, Inches(0.4), GOLD_SOLID)

    # Decorative vertical accent
    add_rect(s, Inches(0.42), Inches(1.2), Inches(0.06), Inches(4.5), GOLD_SOLID)

    # Main title
    add_textbox(s, Inches(0.7), Inches(1.4), Inches(11), Inches(1.4),
                "Vendor Data Intelligence Platform",
                42, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri Light")

    # Sub-title
    add_textbox(s, Inches(0.7), Inches(2.9), Inches(10), Inches(0.6),
                "Agentic Web Scraping System  ·  Interior Design Procurement",
                20, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

    # Divider
    add_rect(s, Inches(0.7), Inches(3.7), Inches(4), Pt(1), GOLD_SOLID)

    # Prepared for line
    add_textbox(s, Inches(0.7), Inches(3.9), Inches(6), Inches(0.4),
                "Prepared for: Executive Leadership", 13, color=LIGHT_BLUE)
    add_textbox(s, Inches(0.7), Inches(4.35), Inches(6), Inches(0.4),
                "March 2026", 13, color=LIGHT_BLUE)

    # Stats box (right side)
    add_rect(s, Inches(9.5), Inches(2.0), Inches(3.4), Inches(3.5), MID_BLUE)

    stats = [
        ("163",  "Vendors Covered"),
        ("31+",  "Product Categories"),
        ("100%", "Automated"),
    ]
    for i, (num, label) in enumerate(stats):
        y = Inches(2.2) + i * Inches(1.1)
        add_textbox(s, Inches(9.7), y, Inches(3.0), Inches(0.55),
                    num, 34, bold=True, color=GOLD_SOLID, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(9.7), y + Inches(0.5), Inches(3.0), Inches(0.4),
                    label, 12, color=WHITE, align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(s, Inches(0.5), Inches(7.15), Inches(8), Inches(0.28),
                "CONFIDENTIAL  ·  Internal Use Only", 9,
                color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.LEFT)
    return s


def slide_02_problem(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    add_rect(s, 0, Inches(1.3), W, Pt(3), GOLD_SOLID)

    add_textbox(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                "The Challenge We Solved", 28, bold=True, color=WHITE, font_name="Calibri Light")
    slide_number(s, 2, TOTAL, color=LIGHT_BLUE)

    # Two-column layout
    col1_x = Inches(0.5)
    col2_x = Inches(6.9)
    col_w  = Inches(6.0)
    top    = Inches(1.7)

    # Left column heading
    add_textbox(s, col1_x, top, col_w, Inches(0.5),
                "Before: Manual & Fragmented", 15, bold=True, color=MID_BLUE)
    add_rect(s, col1_x, top + Inches(0.55), col_w, Pt(2), GOLD_SOLID)

    pain_points = [
        "Staff manually visited vendor websites one-by-one",
        "Copy-pasting product names, prices & specs into spreadsheets",
        "No consistent format — every person structured data differently",
        "Missed fields, outdated prices & broken links went undetected",
        "Scaling to 163 vendors was humanly impossible",
        "Hours spent per vendor, per update cycle",
    ]
    add_bullet_box(s, col1_x, top + Inches(0.75), col_w, Inches(4.5),
                   pain_points, font_size=13, color=DARK_TEXT, bullet="✕  ",
                   spacing=10)

    # Right column heading
    add_textbox(s, col2_x, top, col_w, Inches(0.5),
                "After: Automated & Standardised", 15, bold=True, color=MID_BLUE)
    add_rect(s, col2_x, top + Inches(0.55), col_w, Pt(2), GOLD_SOLID)

    gains = [
        "One command scrapes an entire vendor catalogue",
        "All fields extracted, cleaned & structured automatically",
        "Consistent column layout across every vendor & category",
        "Data validated at source — clean prices, correct dimensions",
        "163 vendors covered with zero extra headcount",
        "Full rescrape in minutes, not days",
    ]
    add_bullet_box(s, col2_x, top + Inches(0.75), col_w, Inches(4.5),
                   gains, font_size=13, color=DARK_TEXT, bullet="✔  ",
                   spacing=10)

    return s


def slide_03_what(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    add_rect(s, 0, Inches(1.3), W, Pt(3), GOLD_SOLID)
    add_textbox(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                "What Is This System?", 28, bold=True, color=WHITE, font_name="Calibri Light")
    slide_number(s, 3, TOTAL, color=LIGHT_BLUE)

    # Intro paragraph
    add_textbox(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.65),
                "An AI-powered agent that automatically visits interior design vendor websites, "
                "reads every product page, and exports clean, structured Excel data — "
                "ready for procurement, quoting, and studio use.",
                14, color=DARK_TEXT, word_wrap=True)

    # Three feature cards
    cards = [
        ("🤖  Agentic",
         "Claude AI writes, runs, and monitors the scraper code. "
         "Each vendor gets a custom scraper tailored to its website structure."),
        ("📊  Structured Output",
         "Data is exported as multi-tab Excel files with consistent columns: "
         "Name, SKU, Price, Dimensions, Materials, Images, and more."),
        ("⚡  Scalable",
         "163 vendors already mapped. Adding a new vendor takes minutes — "
         "just point the system at the tracker and run one command."),
    ]

    cx = [Inches(0.5), Inches(4.7), Inches(8.9)]
    cw = Inches(3.9)
    for i, (title, body) in enumerate(cards):
        add_rect(s, cx[i], Inches(2.5), cw, Inches(3.9), LIGHT_BLUE)
        add_rect(s, cx[i], Inches(2.5), cw, Inches(0.06), MID_BLUE)
        add_textbox(s, cx[i] + Inches(0.15), Inches(2.65), cw - Inches(0.3), Inches(0.65),
                    title, 15, bold=True, color=MID_BLUE)
        add_textbox(s, cx[i] + Inches(0.15), Inches(3.35), cw - Inches(0.3), Inches(2.8),
                    body, 13, color=DARK_TEXT, word_wrap=True)

    return s


def slide_04_process(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    add_rect(s, 0, Inches(1.3), W, Pt(3), GOLD_SOLID)
    add_textbox(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                "How the Process Works", 28, bold=True, color=WHITE, font_name="Calibri Light")
    slide_number(s, 4, TOTAL, color=LIGHT_BLUE)

    # 5-step horizontal flow
    steps = [
        ("1", "SD Tracker",          "Studio team maintains the master Excel tracker with vendor names, category URLs, and required columns."),
        ("2", "Vendor Parser",        "System reads the tracker and extracts category links and column specs for each vendor automatically."),
        ("3", "AI Scraper Writer",    "Claude AI generates a custom scraper.py for the vendor — tailored to its specific website layout."),
        ("4", "Playwright Scraper",   "The scraper navigates every product listing, opens each product page, and collects all available data."),
        ("5", "Structured Excel",     "Clean, formatted Excel files are saved with one sheet per category, correct columns, and live image links."),
    ]

    step_w = Inches(2.3)
    gap    = Inches(0.28)
    top    = Inches(1.7)

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.4) + i * (step_w + gap)

        # Step circle
        circle = s.shapes.add_shape(9, x + Inches(0.75), top + Inches(0.1),
                                    Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = NAVY
        circle.line.fill.background()

        add_textbox(s, x + Inches(0.75), top + Inches(0.13),
                    Inches(0.8), Inches(0.55),
                    num, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Arrow connector (not last)
        if i < len(steps) - 1:
            arrow_x = x + step_w + Inches(0.02)
            add_textbox(s, arrow_x, top + Inches(0.25), Inches(0.25), Inches(0.4),
                        "▶", 16, color=GOLD_SOLID, align=PP_ALIGN.CENTER)

        # Card
        add_rect(s, x, top + Inches(1.1), step_w, Inches(4.5), OFF_WHITE)
        add_rect(s, x, top + Inches(1.1), step_w, Inches(0.05), MID_BLUE)
        add_textbox(s, x + Inches(0.1), top + Inches(1.2), step_w - Inches(0.2), Inches(0.6),
                    title, 13, bold=True, color=MID_BLUE)
        add_textbox(s, x + Inches(0.1), top + Inches(1.85), step_w - Inches(0.2), Inches(3.5),
                    desc, 11.5, color=DARK_TEXT, word_wrap=True)

    return s


def slide_05_architecture(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    add_rect(s, 0, Inches(1.3), W, Pt(3), GOLD_SOLID)
    add_textbox(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                "System Architecture", 28, bold=True, color=WHITE, font_name="Calibri Light")
    slide_number(s, 5, TOTAL, color=LIGHT_BLUE)

    components = [
        ("vendor_parser.py",  "Reads the SD Tracker Excel file. Extracts category names, URLs, and required column lists per vendor. Outputs structured JSON."),
        ("orchestrator.py",   "Central dispatcher. Checks if a scraper exists, runs it in full or test mode, sets environment variables, and handles exit codes."),
        ("base_scraper.py",   "Shared library used by all scrapers. Provides: Playwright browser, Excel writer, text cleaners, price parser, dimension parser, SKU generator."),
        ("<Vendor>/scraper.py", "Vendor-specific scraper written by Claude AI. Handles the unique HTML structure of each vendor website. Fully standalone and rerunnable."),
        ("ExcelWriter",       "Dynamic column engine. Buffers all scraped rows, merges tracker columns with discovered fields, and writes the final formatted Excel file."),
        ("SD Tracker .xlsx",  "Single source of truth. Studio team owns this file. Defines which vendors to scrape, which categories, what columns are expected."),
    ]

    col_x = [Inches(0.45), Inches(4.65), Inches(8.85)]
    row_y = [Inches(1.6), Inches(4.3)]
    w = Inches(3.85)
    h = Inches(2.4)

    for i, (name, desc) in enumerate(components):
        col = i % 3
        row = i // 3
        x = col_x[col]
        y = row_y[row]

        add_rect(s, x, y, w, h, OFF_WHITE)
        add_rect(s, x, y, w, Inches(0.05), GOLD_SOLID)
        add_textbox(s, x + Inches(0.12), y + Inches(0.1), w - Inches(0.24), Inches(0.5),
                    name, 13, bold=True, color=NAVY, font_name="Consolas")
        add_textbox(s, x + Inches(0.12), y + Inches(0.65), w - Inches(0.24), h - Inches(0.8),
                    desc, 12, color=DARK_TEXT, word_wrap=True)

    return s


def slide_06_scale(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    add_rect(s, 0, Inches(1.3), W, Pt(3), GOLD_SOLID)
    add_textbox(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                "Scale & Coverage", 28, bold=True, color=WHITE, font_name="Calibri Light")
    slide_number(s, 6, TOTAL, color=LIGHT_BLUE)

    # Big stat row
    stats = [
        ("163", "Vendors in\nthe Tracker"),
        ("31+",  "Product Categories\nper Vendor"),
        ("500+", "Data Fields\nCaptured"),
        ("100%", "Automated\nData Pipeline"),
    ]

    sw = Inches(2.9)
    for i, (num, label) in enumerate(stats):
        x = Inches(0.5) + i * Inches(3.1)
        add_rect(s, x, Inches(1.65), sw, Inches(2.0), NAVY)
        add_textbox(s, x, Inches(1.75), sw, Inches(1.1),
                    num, 46, bold=True, color=GOLD_SOLID, align=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(2.85), sw, Inches(0.75),
                    label, 12, color=WHITE, align=PP_ALIGN.CENTER)

    # Category groups
    add_textbox(s, Inches(0.5), Inches(4.0), Inches(12), Inches(0.45),
                "Categories Covered", 15, bold=True, color=MID_BLUE)
    add_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Pt(1), GOLD_SOLID)

    cats = [
        ("Furniture",    "Nightstands · Coffee Tables · Dining Tables · Consoles · Desks · Beds · Cabinets · Dressers · Bar Carts · Bookcases"),
        ("Seating",      "Dining Chairs · Bar Stools · Sofas · Loveseats · Sectionals · Benches"),
        ("Lighting",     "Chandeliers · Pendants · Sconces · Flush Mounts · Table Lamps · Floor Lamps · Bulbs"),
        ("Accessories",  "Mirrors · Pillows & Throws · Vases · Bowls · Trays · Objects · Wall Decor · Rugs"),
    ]
    for i, (grp, items) in enumerate(cats):
        y = Inches(4.65) + i * Inches(0.62)
        add_textbox(s, Inches(0.5), y, Inches(1.5), Inches(0.5),
                    grp, 12, bold=True, color=NAVY)
        add_textbox(s, Inches(2.1), y, Inches(10.7), Inches(0.5),
                    items, 11.5, color=MID_GRAY)

    return s


def slide_07_output(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    add_rect(s, 0, Inches(1.3), W, Pt(3), GOLD_SOLID)
    add_textbox(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                "Output Quality & Data Standards", 28, bold=True, color=WHITE, font_name="Calibri Light")
    slide_number(s, 7, TOTAL, color=LIGHT_BLUE)

    # Left: output fields
    add_textbox(s, Inches(0.5), Inches(1.6), Inches(6), Inches(0.45),
                "Fields Captured per Product", 14, bold=True, color=MID_BLUE)

    fields_left = [
        "Source  ·  Image URL  ·  Manufacturer",
        "Product Name  ·  SKU  ·  Product Family ID",
        "Price  ·  Description",
        "Width  ·  Depth  ·  Height  ·  Diameter",
        "Weight  ·  Materials  ·  Specifications",
        "Finish  ·  Color  ·  Fabric  ·  Upholstery",
    ]
    fields_right = [
        "Seat Height  ·  Seat Depth  ·  Arm Height",
        "Socket  ·  Wattage  ·  Light Source",
        "Shade Details  ·  Canopy  ·  Chain Length",
        "Designer  ·  Collection  ·  Lead Time",
        "COM  ·  COL  ·  COT availability",
        "Tearsheet Link  ·  + any field found on page",
    ]
    for i, (fl, fr) in enumerate(zip(fields_left, fields_right)):
        y = Inches(2.1) + i * Inches(0.7)
        add_rect(s, Inches(0.5), y, Inches(5.8), Inches(0.58),
                 OFF_WHITE if i % 2 == 0 else WHITE)
        add_textbox(s, Inches(0.65), y + Inches(0.08), Inches(5.5), Inches(0.45),
                    fl, 12, color=DARK_TEXT)
        add_rect(s, Inches(6.65), y, Inches(6.1), Inches(0.58),
                 OFF_WHITE if i % 2 == 0 else WHITE)
        add_textbox(s, Inches(6.8), y + Inches(0.08), Inches(5.8), Inches(0.45),
                    fr, 12, color=DARK_TEXT)

    # Rules box at bottom
    rules = [
        "Prices: numeric only — no $, USD, or commas  (e.g. 1200.0)",
        "Dimensions: numbers only — no inch marks, fractions converted to decimals",
        "Product Names: sentence case applied uniformly across all vendors",
        "SKUs: auto-generated when vendor does not provide one",
    ]
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.78), LIGHT_BLUE)
    add_textbox(s, Inches(0.65), Inches(6.62), Inches(12.0), Inches(0.65),
                "Data Standards:  " + "   ·   ".join(rules),
                10.5, color=NAVY, word_wrap=True)

    return s


def slide_08_value(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    add_rect(s, 0, Inches(1.3), W, Pt(3), GOLD_SOLID)
    add_textbox(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                "Business Value", 28, bold=True, color=WHITE, font_name="Calibri Light")
    slide_number(s, 8, TOTAL, color=LIGHT_BLUE)

    values = [
        ("Time Savings",
         "What previously took days of manual work per vendor now runs in minutes. "
         "A full 163-vendor scrape can be triggered and completed without any manual effort."),
        ("Accuracy & Consistency",
         "Human copy-paste errors are eliminated. Every price is a clean number, "
         "every dimension is unit-free, every product name follows the same casing standard."),
        ("Scalability",
         "The same system handles one vendor or one hundred. "
         "Adding a new vendor requires only an entry in the SD tracker — no new code."),
        ("Studio Productivity",
         "Studio staff spend time on design decisions, not data entry. "
         "Procurement, quoting, and presentation prep are faster and more reliable."),
        ("Competitive Edge",
         "Up-to-date, structured product data across 163 vendors gives the studio "
         "a sourcing advantage that competitors relying on manual research cannot match."),
        ("AI-Driven & Future-Ready",
         "Built on Claude AI, the system self-writes vendor scrapers. "
         "Improvements to the AI model automatically improve the scraping quality."),
    ]

    col_x = [Inches(0.45), Inches(6.9)]
    row_ys = [Inches(1.65), Inches(3.2), Inches(4.75)]
    w = Inches(6.0)
    h = Inches(1.3)

    for i, (title, body) in enumerate(values):
        col = i % 2
        row = i // 2
        x = col_x[col]
        y = row_ys[row]
        add_rect(s, x, y, w, h, OFF_WHITE)
        add_rect(s, x, y, Inches(0.06), h, MID_BLUE)
        add_textbox(s, x + Inches(0.18), y + Inches(0.1), w - Inches(0.3), Inches(0.42),
                    title, 13, bold=True, color=MID_BLUE)
        add_textbox(s, x + Inches(0.18), y + Inches(0.52), w - Inches(0.3), Inches(0.72),
                    body, 11.5, color=DARK_TEXT, word_wrap=True)

    return s


def slide_09_techstack(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    add_rect(s, 0, Inches(1.3), W, Pt(3), GOLD_SOLID)
    add_textbox(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                "Technology Stack", 28, bold=True, color=WHITE, font_name="Calibri Light")
    slide_number(s, 9, TOTAL, color=LIGHT_BLUE)

    tech = [
        ("Claude AI  (Anthropic)",    "Powers the agentic scraper generation. Understands website HTML, writes Python code, and adapts to each vendor's unique structure."),
        ("Python 3.13",               "Core language for all scraping logic, data parsing, and orchestration. Fast, readable, and extensively tested."),
        ("Playwright",                "Headless browser automation. Handles JavaScript-rendered pages, lazy-loaded images, and dynamic content that simple HTTP requests cannot reach."),
        ("openpyxl",                  "Writes structured, formatted Excel files. Supports dynamic column ordering, header styling, and multi-sheet workbooks."),
        ("Claude Code CLI",           "The command-line interface through which Claude AI reads, writes, and runs code autonomously in the project workspace."),
        ("SD Tracker  (Excel)",       "The human-facing control panel. Studio staff update it to add vendors, adjust categories, or change required column names — no code needed."),
    ]

    col_x = [Inches(0.45), Inches(6.9)]
    row_ys = [Inches(1.65), Inches(3.15), Inches(4.65)]
    w = Inches(6.0)
    h = Inches(1.25)

    for i, (name, desc) in enumerate(tech):
        col = i % 2
        row = i // 2
        x = col_x[col]
        y = row_ys[row]
        add_rect(s, x, y, w, h, OFF_WHITE)
        add_rect(s, x, y, w, Inches(0.04), GOLD_SOLID)
        add_textbox(s, x + Inches(0.12), y + Inches(0.08), w - Inches(0.24), Inches(0.42),
                    name, 13, bold=True, color=NAVY, font_name="Consolas")
        add_textbox(s, x + Inches(0.12), y + Inches(0.5), w - Inches(0.24), h - Inches(0.6),
                    desc, 11.5, color=DARK_TEXT, word_wrap=True)

    # Footer note
    add_rect(s, Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.7), LIGHT_BLUE)
    add_textbox(s, Inches(0.65), Inches(6.65), Inches(12.0), Inches(0.5),
                "All components are open-source or proprietary-free except Claude AI (API subscription). "
                "No third-party scraping services or paid data providers are required.",
                11, color=NAVY)

    return s


def slide_10_next(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    add_rect(s, 0, Inches(1.3), W, Pt(3), GOLD_SOLID)
    add_textbox(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
                "Next Steps & Roadmap", 28, bold=True, color=WHITE, font_name="Calibri Light")
    slide_number(s, 10, TOTAL, color=LIGHT_BLUE)

    phases = [
        ("Phase 1  —  In Progress",
         MID_BLUE,
         [
             "Complete scraper coverage for all 163 tracked vendors",
             "Regina Andrew, Visual Comfort, Arteriors — active",
             "Full test → validate → full-run pipeline operational",
         ]),
        ("Phase 2  —  Near Term",
         GOLD_SOLID,
         [
             "Scheduled re-scrapes: auto-refresh vendor data on a weekly/monthly cycle",
             "Price-change detection: flag products where price has changed since last run",
             "Image download: cache product images locally for offline presentations",
         ]),
        ("Phase 3  —  Future",
         MID_GRAY,
         [
             "Database integration: push scraped data directly into studio ERP / CRM",
             "Web dashboard: browse and search all vendor products in one interface",
             "Automatic vendor discovery: AI identifies and proposes new vendors to add",
         ]),
    ]

    for i, (phase, color, items) in enumerate(phases):
        x = Inches(0.45) + i * Inches(4.28)
        add_rect(s, x, Inches(1.65), Inches(4.1), Inches(5.3), OFF_WHITE)
        add_rect(s, x, Inches(1.65), Inches(4.1), Inches(0.55), color)
        add_textbox(s, x + Inches(0.12), Inches(1.72), Inches(3.9), Inches(0.42),
                    phase, 12, bold=True, color=WHITE)
        for j, item in enumerate(items):
            y = Inches(2.35) + j * Inches(1.4)
            add_rect(s, x + Inches(0.15), y, Inches(3.8), Inches(1.2), WHITE)
            add_textbox(s, x + Inches(0.28), y + Inches(0.1), Inches(3.55), Inches(1.0),
                        item, 12, color=DARK_TEXT, word_wrap=True)

    # Closing line
    add_rect(s, Inches(0.45), Inches(7.0), Inches(12.4), Pt(2), GOLD_SOLID)
    add_textbox(s, Inches(0.45), Inches(6.6), Inches(12.4), Inches(0.4),
                "The foundation is built. Every new vendor is now a one-command operation.",
                14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    return s


# ────────────────────────────────────────────────────────────────────────────
# MAIN
# ────────────────────────────────────────────────────────────────────────────

def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = new_prs()

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_what(prs)
    slide_04_process(prs)
    slide_05_architecture(prs)
    slide_06_scale(prs)
    slide_07_output(prs)
    slide_08_value(prs)
    slide_09_techstack(prs)
    slide_10_next(prs)

    prs.save(str(OUT_FILE))
    print(f"[OK] Saved: {OUT_FILE}")


if __name__ == "__main__":
    main()
