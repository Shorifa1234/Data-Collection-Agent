from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette ──────────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x23, 0x3E)
TEAL      = RGBColor(0x00, 0x9B, 0x94)
GOLD      = RGBColor(0xD4, 0xAF, 0x37)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF9)
MID_GREY  = RGBColor(0x6B, 0x7A, 0x99)
DARK_TEXT = RGBColor(0x1A, 0x23, 0x3E)
GREEN     = RGBColor(0x26, 0xA6, 0x9A)
PURPLE    = RGBColor(0x5C, 0x6B, 0xC0)
ORANGE    = RGBColor(0xEF, 0x6C, 0x00)
DARK_GREEN= RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN=RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEB)
RED_TEXT  = RGBColor(0x7B, 0x1F, 0x1F)
RED_BAR   = RGBColor(0xC0, 0x39, 0x39)
GREEN_TEXT= RGBColor(0x1B, 0x5E, 0x20)
GREEN_BORD= RGBColor(0x66, 0xBB, 0x6A)
CARD_BG   = RGBColor(0x24, 0x2F, 0x55)
TEAL_DIM  = RGBColor(0x00, 0x6B, 0x66)
TEAL_LIGHT= RGBColor(0xA0, 0xC4, 0xC2)
CARD_BDR  = RGBColor(0xDD, 0xE3, 0xED)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Helper utilities ───────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    f = shape.fill
    if fill_rgb:
        f.solid(); f.fore_color.rgb = fill_rgb
    else:
        f.background()
    ln = shape.line
    if line_rgb:
        ln.color.rgb = line_rgb
        ln.width = Pt(line_w)
    else:
        ln.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = "Calibri"
    return txb

def full_bg(slide, color):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=color)

def accent_bar(slide, color=TEAL):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_rgb=color)

def bottom_bar(slide):
    add_rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55), fill_rgb=NAVY)

def slide_num(slide, n):
    add_text(slide, str(n),
             SLIDE_W - Inches(0.6), SLIDE_H - Inches(0.45),
             Inches(0.4), Inches(0.35),
             font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, NAVY)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill_rgb=TEAL)
add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill_rgb=GOLD)
add_rect(slide, Inches(8.8), Inches(0.8), Inches(4.5), Inches(4.5), fill_rgb=TEAL_DIM)

add_text(slide, "Meet Your New\nDigital Assistant",
         Inches(0.7), Inches(1.5), Inches(7.5), Inches(2.8),
         font_size=52, bold=True, color=WHITE)
add_text(slide, "How our smart tool saves hours of work every week\nand keeps our product data always up to date",
         Inches(0.7), Inches(4.4), Inches(7.0), Inches(1.6),
         font_size=20, color=TEAL_LIGHT)
add_rect(slide, Inches(0.7), Inches(6.0), Inches(3.6), Inches(0.55), fill_rgb=TEAL)
add_text(slide, "  Smarter. Faster. Always Accurate.",
         Inches(0.7), Inches(6.0), Inches(3.6), Inches(0.55),
         font_size=15, bold=True, color=WHITE)
bottom_bar(slide)
slide_num(slide, 1)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 – THE OLD WAY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, LIGHT_BG)
accent_bar(slide, TEAL)

add_text(slide, "The Old Way — Lots of Manual Work",
         Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
         font_size=30, bold=True, color=NAVY)
add_text(slide, "Every time a vendor updates prices, images or specs, someone had to do this by hand:",
         Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.6),
         font_size=17, color=MID_GREY)

cards = [
    ("Hours of Browsing",  "Visiting each vendor website one by one"),
    ("Copy & Paste",       "Manually copying names, prices and specs into spreadsheets"),
    ("Human Errors",       "Typos, missed products, outdated info sneaking in"),
    ("Never Ending",       "Vendors update constantly — the work never stops"),
]
bar_colors = [TEAL, GOLD, RED_BAR, NAVY]
for i, (title, desc) in enumerate(cards):
    x = Inches(0.45) + i * Inches(3.2)
    add_rect(slide, x, Inches(1.9), Inches(3.0), Inches(3.6),
             fill_rgb=WHITE, line_rgb=CARD_BDR, line_w=1)
    add_rect(slide, x, Inches(1.9), Inches(3.0), Inches(0.09), fill_rgb=bar_colors[i])
    add_text(slide, title, x + Inches(0.12), Inches(2.1), Inches(2.8), Inches(0.5),
             font_size=16, bold=True, color=NAVY)
    add_text(slide, desc, x + Inches(0.12), Inches(2.75), Inches(2.8), Inches(1.8),
             font_size=13, color=MID_GREY)

bottom_bar(slide)
slide_num(slide, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 – THE SOLUTION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, NAVY)
accent_bar(slide, GOLD)

add_text(slide, "Introducing Our\nSmart Scraping Tool",
         Inches(0.7), Inches(0.7), Inches(7), Inches(1.8),
         font_size=38, bold=True, color=WHITE)
add_text(slide, "Think of it as a tireless team member who visits every vendor\nwebsite automatically and fills in the spreadsheet for you — perfectly, every time.",
         Inches(0.7), Inches(2.6), Inches(7.2), Inches(1.4),
         font_size=18, color=TEAL_LIGHT)

steps = [
    ("1", "We point it\nat a vendor",     TEAL),
    ("2", "It reads every\nproduct page",  GOLD),
    ("3", "Ready-to-use\nspreadsheet",     RGBColor(0x4C, 0xAF, 0x50)),
]
for i, (num, label, col) in enumerate(steps):
    x = Inches(0.8) + i * Inches(3.8)
    add_rect(slide, x, Inches(4.2), Inches(2.8), Inches(2.4), fill_rgb=col)
    add_text(slide, num,   x + Inches(0.15), Inches(4.3),  Inches(0.6), Inches(0.7),
             font_size=36, bold=True, color=WHITE)
    add_text(slide, label, x + Inches(0.15), Inches(5.1),  Inches(2.5), Inches(1.2),
             font_size=17, bold=True, color=WHITE)
    if i < 2:
        add_text(slide, "->", x + Inches(2.85), Inches(4.85), Inches(0.7), Inches(0.7),
                 font_size=28, bold=True, color=WHITE)

bottom_bar(slide)
slide_num(slide, 3)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 – WHAT IT COLLECTS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, LIGHT_BG)
accent_bar(slide, TEAL)

add_text(slide, "What Information Does It Gather?",
         Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
         font_size=30, bold=True, color=NAVY)
add_text(slide, "For every single product on a vendor's website, it automatically captures:",
         Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         font_size=17, color=MID_GREY)

items = [
    ("Product Photos",    "Main image and all variant photos"),
    ("Pricing",           "Always the latest list price — no guessing"),
    ("Dimensions",        "Height, width, depth, weight — all units clean"),
    ("Full Description",  "Product story, materials, and specifications"),
    ("Variants & Finishes","Every colour, finish or size option listed"),
    ("Product Codes",     "SKU codes generated if the vendor doesn't provide them"),
    ("Designer & Origin", "Who made it, where it comes from"),
    ("Tearsheet Links",   "Direct link to the product tearsheet if available"),
]
col1, col2 = items[:4], items[4:]
for col_items, cx in [(col1, Inches(0.4)), (col2, Inches(6.8))]:
    for j, (title, desc) in enumerate(col_items):
        y = Inches(1.85) + j * Inches(1.22)
        add_rect(slide, cx, y, Inches(6.0), Inches(1.05),
                 fill_rgb=WHITE, line_rgb=CARD_BDR, line_w=1)
        add_text(slide, title, cx + Inches(0.2), y + Inches(0.06), Inches(5.6), Inches(0.42),
                 font_size=15, bold=True, color=NAVY)
        add_text(slide, desc,  cx + Inches(0.2), y + Inches(0.5),  Inches(5.6), Inches(0.48),
                 font_size=12, color=MID_GREY)

bottom_bar(slide)
slide_num(slide, 4)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 – KEY BENEFITS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, WHITE)
accent_bar(slide, GOLD)

add_text(slide, "Why Your Team Will Love This",
         Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
         font_size=30, bold=True, color=NAVY)

benefits = [
    (TEAL,   "Huge Time Savings",    "What used to take days now\ntakes minutes automatically"),
    (GOLD,   "Always Accurate",      "No typos, no missing fields,\nno outdated prices"),
    (PURPLE, "Organised & Ready",    "Data lands straight into a\nneat formatted spreadsheet"),
    (GREEN,  "Stays Up to Date",     "Run it any time — vendor\nupdates captured instantly"),
    (ORANGE, "Works for Every Vendor","One tool covers all vendors\nwe work with"),
    (RGBColor(0x8E, 0x24, 0xAA), "No Manual Checking",
     "Your team can focus on design,\nnot data entry"),
]
for i, (col, title, desc) in enumerate(benefits):
    row, col_n = divmod(i, 3)
    x = Inches(0.4) + col_n * Inches(4.28)
    y = Inches(1.1)  + row   * Inches(2.8)
    add_rect(slide, x, y, Inches(3.9), Inches(2.55), fill_rgb=col)
    add_text(slide, title, x + Inches(0.2), y + Inches(0.25), Inches(3.5), Inches(0.6),
             font_size=17, bold=True, color=WHITE)
    add_text(slide, desc,  x + Inches(0.2), y + Inches(1.0),  Inches(3.5), Inches(1.3),
             font_size=13, color=WHITE)

bottom_bar(slide)
slide_num(slide, 5)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 – BEFORE vs AFTER
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, LIGHT_BG)
accent_bar(slide, TEAL)

add_text(slide, "Before & After",
         Inches(0.5), Inches(0.3), Inches(10), Inches(0.65),
         font_size=30, bold=True, color=NAVY)

# Left panel
add_rect(slide, Inches(0.4), Inches(1.1), Inches(5.8), Inches(5.8),
         fill_rgb=LIGHT_RED, line_rgb=RGBColor(0xE0, 0x80, 0x80), line_w=1)
add_rect(slide, Inches(0.4), Inches(1.1), Inches(5.8), Inches(0.55), fill_rgb=RED_BAR)
add_text(slide, "Before", Inches(0.55), Inches(1.13), Inches(5.5), Inches(0.5),
         font_size=18, bold=True, color=WHITE)

before_items = [
    "Open vendor website manually",
    "Search through hundreds of products",
    "Copy each product name, price, dimensions",
    "Paste into spreadsheet, fix formatting",
    "Check for typos and missing info",
    "Repeat for every vendor, every time",
    "Hope nothing changed since last time",
]
for j, item in enumerate(before_items):
    add_text(slide, "  X  " + item,
             Inches(0.6), Inches(1.85) + j * Inches(0.63), Inches(5.4), Inches(0.55),
             font_size=13, color=RED_TEXT)

# Right panel
add_rect(slide, Inches(7.0), Inches(1.1), Inches(5.8), Inches(5.8),
         fill_rgb=LIGHT_GREEN, line_rgb=GREEN_BORD, line_w=1)
add_rect(slide, Inches(7.0), Inches(1.1), Inches(5.8), Inches(0.55), fill_rgb=DARK_GREEN)
add_text(slide, "After", Inches(7.15), Inches(1.13), Inches(5.5), Inches(0.5),
         font_size=18, bold=True, color=WHITE)

after_items = [
    "Type the vendor name",
    "Press Enter",
    "Go for a coffee",
    "Come back to a perfect spreadsheet",
    "Every product — captured",
    "Every field — filled in correctly",
    "Ready to use immediately",
]
for j, item in enumerate(after_items):
    add_text(slide, "  OK  " + item,
             Inches(7.2), Inches(1.85) + j * Inches(0.63), Inches(5.4), Inches(0.55),
             font_size=13, color=GREEN_TEXT)

bottom_bar(slide)
slide_num(slide, 6)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 – WHO BENEFITS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, NAVY)
accent_bar(slide, GOLD)

add_text(slide, "Who Benefits in Our Office?",
         Inches(0.7), Inches(0.4), Inches(9), Inches(0.7),
         font_size=30, bold=True, color=WHITE)
add_text(slide, "This tool touches every team that works with product information:",
         Inches(0.7), Inches(1.15), Inches(11), Inches(0.5),
         font_size=17, color=TEAL_LIGHT)

roles = [
    ("Design Teams",      "Instantly browse the full product range\nof any vendor — no more hunting"),
    ("Procurement",       "Accurate prices and specs, always\nup to date — order with confidence"),
    ("Project Managers",  "Products ready for client proposals\nwithout waiting for data entry"),
    ("Marketing",         "Fresh images and descriptions\nfor every campaign, automatically"),
    ("Studio Operations", "Clean, consistent spreadsheets —\nno reformatting, no chasing people"),
]
for i, (role, desc) in enumerate(roles):
    x = Inches(0.5) + (i % 3) * Inches(4.2)
    y = Inches(1.9)  + (i // 3) * Inches(2.5)
    add_rect(slide, x, y, Inches(3.8), Inches(2.1),
             fill_rgb=CARD_BG, line_rgb=TEAL, line_w=1)
    add_text(slide, role, x + Inches(0.2), y + Inches(0.12), Inches(3.4), Inches(0.55),
             font_size=15, bold=True, color=GOLD)
    add_text(slide, desc, x + Inches(0.2), y + Inches(0.75), Inches(3.4), Inches(1.2),
             font_size=12, color=WHITE)

bottom_bar(slide)
slide_num(slide, 7)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 – HOW TO USE IT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, WHITE)
accent_bar(slide, TEAL)

add_text(slide, "So... How Does Someone Actually Use It?",
         Inches(0.5), Inches(0.3), Inches(11), Inches(0.7),
         font_size=28, bold=True, color=NAVY)
add_text(slide, "You don't need any technical knowledge. It's as simple as:",
         Inches(0.5), Inches(1.05), Inches(11), Inches(0.5),
         font_size=17, color=MID_GREY)

steps_detail = [
    (TEAL,   "Step 1", "Tell it the vendor name",
     "Just type the vendor name\nand press Enter. That is it."),
    (GOLD,   "Step 2", "It does a quick test first",
     "It checks a few products to\nmake sure everything looks right."),
    (PURPLE, "Step 3", "Confirm and go",
     "One confirmation and it runs\nthe full catalogue automatically."),
    (GREEN,  "Step 4", "Pick up your file",
     "Find the finished Excel file\nin the vendor folder. Done!"),
]
for i, (col, step, title, desc) in enumerate(steps_detail):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.75)
    add_rect(slide, x + Inches(0.9), y, Inches(1.05), Inches(1.05), fill_rgb=col)
    add_text(slide, str(i + 1), x + Inches(0.9), y, Inches(1.05), Inches(1.05),
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, step,  x + Inches(0.2), y + Inches(1.2), Inches(2.8), Inches(0.45),
             font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.2), y + Inches(1.65), Inches(2.8), Inches(0.55),
             font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, desc,  x + Inches(0.2), y + Inches(2.3),  Inches(2.8), Inches(1.2),
             font_size=13, color=MID_GREY, align=PP_ALIGN.CENTER)

add_text(slide, "No coding. No setup. No training needed.",
         Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5),
         font_size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

bottom_bar(slide)
slide_num(slide, 8)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 – REAL IMPACT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, LIGHT_BG)
accent_bar(slide, TEAL)

add_text(slide, "What This Means in Real Life",
         Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
         font_size=30, bold=True, color=NAVY)

stats = [
    (TEAL,       "100s",    "of products captured\nper vendor, automatically"),
    (NAVY,       "Minutes", "to collect what used\nto take days of manual work"),
    (GOLD,       "Zero",    "copy-paste errors —\ndata captured directly"),
    (DARK_GREEN, "Always",  "up to date whenever\nyou need it"),
]
for i, (col, big, small) in enumerate(stats):
    x = Inches(0.4) + i * Inches(3.2)
    add_rect(slide, x, Inches(1.2), Inches(2.9), Inches(3.2), fill_rgb=col)
    add_text(slide, big,   x + Inches(0.15), Inches(1.5),  Inches(2.6), Inches(1.2),
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, small, x + Inches(0.15), Inches(2.85), Inches(2.6), Inches(1.3),
             font_size=15, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Our team spends less time on data entry\nand more time on the work that actually matters.",
         Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.0),
         font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(3.5), Inches(5.95), Inches(6.3), Inches(0.65), fill_rgb=TEAL)
add_text(slide, "More creativity. Less spreadsheet headaches.",
         Inches(3.5), Inches(5.95), Inches(6.3), Inches(0.65),
         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

bottom_bar(slide)
slide_num(slide, 9)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 – CLOSING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
full_bg(slide, NAVY)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill_rgb=GOLD)
add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill_rgb=TEAL)

add_text(slide, "We're Already Live.",
         Inches(0.7), Inches(0.7), Inches(9), Inches(1.1),
         font_size=50, bold=True, color=WHITE)
add_text(slide, "The tool is running today and delivering results for our team.",
         Inches(0.7), Inches(1.85), Inches(10.5), Inches(0.6),
         font_size=20, color=TEAL_LIGHT)

whats_next = [
    ("Currently active",   "Running for multiple vendors right now"),
    ("More vendors coming", "Easy to add — just tell us which one"),
    ("Open to your ideas", "What other manual tasks can we automate next?"),
]
for i, (title, desc) in enumerate(whats_next):
    y = Inches(2.7) + i * Inches(1.1)
    add_rect(slide, Inches(0.7), y, Inches(11.5), Inches(0.9), fill_rgb=CARD_BG)
    add_text(slide, title, Inches(0.9),  y + Inches(0.08), Inches(4.0), Inches(0.55),
             font_size=16, bold=True, color=GOLD)
    add_text(slide, desc,  Inches(5.2),  y + Inches(0.08), Inches(7.0), Inches(0.55),
             font_size=15, color=WHITE)

add_text(slide, "Questions? We would love to show you a live demo!",
         Inches(0.7), Inches(6.15), Inches(11.5), Inches(0.55),
         font_size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

slide_num(slide, 10)

# ── Save ──────────────────────────────────────────────────────
out = r"d:\workspace\playyeard\Agent\Agentic system SKU\Project Overview\Smart Scraping Tool - Office Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
